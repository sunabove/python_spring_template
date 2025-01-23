import argparse
import os
import re

# 환경 변수 설정
os.environ['DOTNET_SYSTEM_GLOBALIZATION_INVARIANT'] = 'true'

def import_or_install(package):
    import importlib
    import subprocess
    import sys

    try:
        # 패키지를 임포트 시도
        return importlib.import_module( package )
    except ImportError:
        # 패키지가 없으면 설치
        print(f"'{package}' 패키지가 없으므로 설치를 진행합니다...")
        subprocess.check_call([sys.executable, "-m", "pip", "install", package])
        return importlib.import_module(package)
    pass
pass

import_or_install( "aspose.slides" ) # Aspose.Slides 패키지 설치

import aspose.slides as slides
from pathlib import Path
import json

from pptx import Presentation

def extract_slide_titles_to_json(ppt_file_path, output_json_path):
    print()
    print(f"Extracting slide titles from {ppt_file_path}...")

    ppt_file_path = Path(ppt_file_path)
    output_json_path = Path(output_json_path)

    # PowerPoint 프레젠테이션 열기
    presentation = Presentation(ppt_file_path)
    slide_titles = []

    for idx, slide in enumerate( presentation.slides ):
        title = None
        slide_number = idx + 1

        # 슬라이드의 제목을 찾아 추출
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame is not None:
                t = shape.text_frame.text.strip()
                if t and "" not in t :  # 텍스트가 비어있지 않은지 확인
                    title = shape.text_frame.text.strip()
                    # 정규 표현식을 사용하여 모든 화이트스페이스 문자를 단일 스페이스로 대체
                    title = re.sub(r'[\u25A0-\u25FF\u2B50\uFFFD\u003F]+', '', title).strip()
                    title = re.sub(r'\s+', ' ', title).strip()
                    title = title[ 0 : 24 ]
                    break
                pass
            pass
        pass

        # 슬라이드 제목 저장
        slide_titles.append({
            "number": slide_number,
            "title": title if title else ""
        })
    pass

    # JSON 파일로 저장
    with open(output_json_path, 'w', encoding='utf-8') as json_file:
        json.dump(slide_titles, json_file, indent=4, ensure_ascii=False)
    pass

    print(f"Slide titles have been saved to {output_json_path}")
    print()
pass # extract_slide_titles_to_json

def convert_ppt_to_html( ppt_file ):

    print(f"Converting ppt file to html : {ppt_file} ...", flush=1 )

    # Load the presentation file
    ppt = slides.Presentation( str(ppt_file) )

    # Create HTML options
    options = slides.export.HtmlOptions()
    # Create a responsive HTML controller
    controller = slides.export.ResponsiveHtmlController() 
    # Set controller as HTML formatter
    options.html_formatter = slides.export.HtmlFormatter.create_custom_formatter(controller)

    folder_path = Path( ppt_file ).parent

    # 모든 HTML 파일을 삭제합니다.
    print( f"deleting all html files in the folder({folder_path})..." )
    [ html_file.unlink() for html_file in folder_path.glob("*.html") ]
    
    # 각 슬라이드를 개별 HTML 파일로 변환합니다.
    for i, slide in enumerate( ppt.slides ):
        # 새로운 프레젠테이션 객체를 생성합니다.
        single_slide_presentation = slides.Presentation()
        # 모든 슬라이드를 제거합니다.
        single_slide_presentation.slides.remove_at(0)
        # 현재 슬라이드를 추가합니다.
        single_slide_presentation.slides.insert_clone(0, slide) 
        
        # 슬라이드를 HTML로 저장합니다.
        html_file_path = str(folder_path / f"{i+1:03d}.html")
        
        print( f"saving a ppt slide to {html_file_path} ..." )    
        single_slide_presentation.save( html_file_path, slides.export.SaveFormat.HTML, options )

        content = ""
        
        # convert utf-8-bom to utf-8
        with open(html_file_path, 'r', encoding='utf-8-sig') as file:
            # 파일의 모든 내용을 읽어들입니다.
            content = file.read().strip()
        pass

        with open(html_file_path, 'w', encoding='utf-8') as file:
            file.write( content )
        pass
        # // convert utf-8-bom to utf-8

        #continue 

        with open(html_file_path, 'r', encoding='utf-8') as file:
            # 파일의 모든 내용을 읽어들입니다.
            content = file.read().strip()
        pass

        with open(html_file_path, 'w', encoding='utf-8') as file:

            if "Created with Aspose.Slides" in content:
                # remove water-mark
                idxApose = content.index( "Created with Aspose.Slides" )
                idxGtransform = content.rfind( "<g transform=\"matrix", 0, idxApose )

                idxPtyLtd = content.rindex( "Aspose Pty Ltd.", idxApose )
                idxG = content.index( "</g>", idxPtyLtd )

                a = content[ 0 : idxGtransform ]
                b = content[ idxG + 4 : ]
                content = a + b
                # // remove water-mark
            pass

            if True : 
                # 슬라이드 제목 제거
                slide_title_tag = "<div class=\"slideTitle\">"
                slide_title_end_tag = "</div>"
                
                idxSlideTitleStart = content.find(slide_title_tag)
                if idxSlideTitleStart > -1:
                    idxSlideTitleClose = content.find(slide_title_end_tag, idxSlideTitleStart)
                    if idxSlideTitleClose > -1:
                        a = content[0:idxSlideTitleStart]
                        b = content[idxSlideTitleClose + len(slide_title_end_tag):]
                        content = a + b
                    pass
                pass
            pass

            if True :
                # 1. viewBox 수정
                content = content.replace('viewBox="0 0 720 540"', 'viewBox="0 0 960 540"')

                # 2. <image> 태그의 width와 height 수정
                content = content.replace('width="720"', 'width="960"')
                content = content.replace('height="540"', 'height="540"')

                # 2. CSS 수정
                new_css = """
                .slides-canvas {
                    display: block;
                    margin: 0 auto;
                    width: 100%;
                    height: auto;
                    aspect-ratio: 16 / 9;
                }
                """
                content = content.replace( '<style>.slide {', '<style>' + new_css + '}\n.slide {' )
            pass

            file.write( content )
        pass
    pass

    print( f"Converted {len(ppt.slides)} slides to HTML files." )
pass # convert_slide_to_html

def check_and_convert_ppts_in_folder(folder_path):
    debug = False 

    print( f"Checking and converting PPT files in folder {folder_path}..." )

    # Recursively search for .pptx files
    for ppt_file in Path(folder_path).rglob("*.pptx"):
        # Get the modification time of the PPT file
        ppt_mod_time = ppt_file.stat().st_mtime

        # Get the list of all HTML files for the corresponding PPT file
        html_files = [ ( ppt_file.parent / f ) for f in os.listdir(ppt_file.parent) if f.endswith('.html')]

        print( f"html_files len: {len( html_files )}" )

        # Check if the PPT file's modification time is newer than all HTML files
        ppt_needs_conversion = False
        
        if len( html_files ) == 0 :
            ppt_needs_conversion = True
        elif len( html_files ) < len( slides.Presentation( str(ppt_file) ).slides ) : 
            ppt_needs_conversion = True
        pass

        if not ppt_needs_conversion :
            for html_file in html_files :
                html_mod_time = html_file.stat().st_mtime 
                
                debug and print( f"ppt_mod_time = {ppt_mod_time}, html_mod_time = {html_mod_time}" )
                
                if ppt_mod_time > html_mod_time:
                    ppt_needs_conversion = True
                    break
                pass
            pass
        pass

        print( f"ppt_needs_conversion = {ppt_needs_conversion}" )

        json_file = ppt_file.parent / "slides.json"

        if ppt_needs_conversion:
            print(f"PPT file {ppt_file} is newer than corresponding HTML files. Converting all slides...")
            # Convert all slides to HTML
            convert_ppt_to_html( ppt_file )

            extract_slide_titles_to_json( ppt_file, json_file )
        else:
            print(f"All HTML files are up-to-date. Skipped converting PPT file {ppt_file}.")
        pass

        json_file = ppt_file.parent / "slides.json" 

        if True or not json_file.exists() :
            extract_slide_titles_to_json( ppt_file, json_file )
        pass
    pass
pass # check_and_convert_ppts_in_folder

if __name__ == "__main__":

    target_folder = "200_ppt_docs"
    current_dir = Path(__file__).parent
    # 현재 디렉토리와 상위 디렉토리에서 폴더 검색
    search_folder = next((f for f in current_dir.glob("../**/") if f.is_dir() and f.name == target_folder), None)

    print( "search_folder = " , search_folder )
    
    parser = argparse.ArgumentParser()
   
    # parser.add_argument( "folder",  help="Folder path to search for ppt files" ) 
    parser.add_argument(
        "folder",
        help="Folder path to search for ppt files",
        nargs="?",  # 위치 인자를 선택적으로 변경
        default=search_folder  # 기본값 설정
    )

    args = parser.parse_args()

    check_and_convert_ppts_in_folder(args.folder)
pass